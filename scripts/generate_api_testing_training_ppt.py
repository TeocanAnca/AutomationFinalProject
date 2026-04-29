import os
from pptx import Presentation
from pptx.util import Inches, Pt


def add_bullet_slide(prs, title, bullets, notes=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    body = slide.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)

    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes


def main():
    prs = Presentation()

    # Improve readability on common 16:9 projectors
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    add_bullet_slide(
        prs,
        "API Testing Fundamentals (Why It Matters)",
        [
            "API tests validate contracts, behavior, reliability, and error handling.",
            "They run faster than UI tests and give earlier feedback in CI.",
            "Core HTTP concepts: methods, status codes, headers, request/response body.",
            "Cover positive, negative, and edge-case scenarios.",
            "Automation brings repeatability, confidence, and scalable coverage.",
        ],
        notes=(
            "Use one real request/response example to explain terms. "
            "Position API tests between unit and UI tests in the test pyramid."
        ),
    )

    # Slide 2
    add_bullet_slide(
        prs,
        "RestAssured Basics in Java",
        [
            "RestAssured provides a fluent syntax: given() -> when() -> then().",
            "Set base URI, headers, and content type to start quickly.",
            "Integrates well with TestNG, Hamcrest, and Jackson.",
            "Supports CRUD flows clearly: POST, GET, PUT, DELETE.",
            "Readable tests help manual testers transition into automation.",
        ],
        notes=(
            "Show a small example using given/when/then to highlight readability."
        ),
    )

    # Slide 3
    add_bullet_slide(
        prs,
        "RequestSpecification: Reuse and Consistency",
        [
            "RequestSpecification centralizes common request setup.",
            "Define base URI, headers, auth, and SSL config once.",
            "Keep tests focused on endpoint behavior, not plumbing.",
            "Reduces duplication and simplifies maintenance across suites.",
            "In this framework: BaseConfig + BaseTest build shared setup.",
        ],
        notes=(
            "Explain \"change once, apply everywhere\" with a header or URL update example."
        ),
    )

    # Slide 4
    add_bullet_slide(
        prs,
        "DTOs/POJOs for Clean Test Design",
        [
            "DTOs model request/response payloads as Java objects.",
            "Avoid fragile hardcoded JSON strings in test methods.",
            "Compile-time field names improve safety and readability.",
            "Test data factories create reusable, consistent object instances.",
            "In this framework: PetDto + TestDataFactory.",
        ],
        notes=(
            "Compare a raw JSON string test with a POJO-based test to show clarity gains."
        ),
    )

    # Slide 5
    add_bullet_slide(
        prs,
        "Serialization and Deserialization (Jackson)",
        [
            "Serialization: POJO -> JSON for outbound request bodies.",
            "Deserialization: JSON -> POJO/List for response validation.",
            "RestAssured can auto-integrate with Jackson for both directions.",
            "Ignore unknown fields to keep tests resilient to non-breaking API changes.",
            "Object-based assertions are cleaner than deep JsonPath strings.",
        ],
        notes=(
            "Mention common pitfalls: field name mismatch, type mismatch, null handling."
        ),
    )

    # Slide 6
    add_bullet_slide(
        prs,
        "Assertions and Validation Strategy",
        [
            "Validate status code, headers, and body in each critical flow.",
            "Use ResponseSpecification for shared response expectations.",
            "Prefer expressive Hamcrest matchers for clear failures.",
            "Include negative checks (e.g., GET after DELETE returns 404).",
            "Use DataProvider-driven tests to widen coverage with less code.",
        ],
        notes=(
            "Show how one data-driven test can validate multiple statuses quickly."
        ),
    )

    # Slide 7
    add_bullet_slide(
        prs,
        "Framework Best Practices for Stable API Automation",
        [
            "Keep clear layers: config, base test, dto, utils, and test classes.",
            "Generate unique test data and always clean up created resources.",
            "Separate environment settings from test logic.",
            "Add request/response logging for fast troubleshooting.",
            "Optimize for CI: deterministic, isolated, and maintainable tests.",
        ],
        notes=(
            "Close with a maturity path: start basic, then refactor into reusable architecture."
        ),
    )

    output = "/Users/manuel.negrean/Downloads/API_automation_training_framework/docs/API_Testing_RestAssured_Training.pptx"
    os.makedirs(os.path.dirname(output), exist_ok=True)
    prs.save(output)
    print(output)


if __name__ == "__main__":
    main()
